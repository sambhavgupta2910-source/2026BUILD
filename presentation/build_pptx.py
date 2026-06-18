#!/usr/bin/env python3
"""
Build the Elevate Homes end-to-end investor + feasibility deck as an editable
.pptx (brand-styled, real tables, speaker notes). Pure python-pptx.

    python3 presentation/build_pptx.py
    -> presentation/Elevate-Homes-End-to-End-Deck.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG   = RGBColor(0x0C, 0x11, 0x18)
SURF = RGBColor(0x16, 0x1F, 0x2B)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
GOLD2= RGBColor(0xE8, 0xD5, 0xA3)
TEXT = RGBColor(0xF0, 0xEC, 0xE4)
SEC  = RGBColor(0xC2, 0xB8, 0x9A)
MUTED= RGBColor(0x7A, 0x6F, 0x5E)
GREEN= RGBColor(0x34, 0xD3, 0x99)
RED  = RGBColor(0xC0, 0x55, 0x45)
DISP = "Cormorant Garamond"
BODY = "Manrope"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def para(tf, text, size, color, bold=False, italic=False, font=BODY, align=PP_ALIGN.LEFT,
         space_after=6, first=False, line=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font
    f.color.rgb = color
    return p


def panel(s, l, t, w, h, fill=SURF, line=GOLD):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def kicker_title(s, kicker, title, tsize=34):
    _, tf = box(s, 0.7, 0.55, 12, 1.7)
    para(tf, kicker.upper(), 12, GOLD, bold=True, font=BODY, space_after=6, first=True)
    for i, line in enumerate(title.split("\n")):
        para(tf, line, tsize, GOLD2, bold=True, font=DISP, space_after=2, line=1.0)


def card(s, l, t, w, h, heading, lines, hcolor=GOLD2, accent=GOLD):
    panel(s, l, t, w, h, fill=SURF, line=accent)
    _, tf = box(s, l + 0.25, t + 0.18, w - 0.5, h - 0.36)
    para(tf, heading, 14, hcolor, bold=True, font=BODY, space_after=6, first=True)
    for ln in lines:
        para(tf, ln, 11.5, SEC, font=BODY, space_after=4, line=1.05)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def table(s, l, t, w, rows, col_w, header_fill=RGBColor(0x2A, 0x24, 0x12), row_h=0.42, fsize=11):
    nrows, ncols = len(rows), len(rows[0])
    gf = s.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(row_h * nrows))
    tbl = gf.table
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if ri == 0 else (RGBColor(0x12, 0x18, 0x20) if ri % 2 else RGBColor(0x0F, 0x14, 0x1B))
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tfc = cell.text_frame; tfc.word_wrap = True
            p = tfc.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fsize)
            r.font.name = BODY
            r.font.bold = (ri == 0)
            r.font.color.rgb = GOLD if ri == 0 else (TEXT if ci == 0 else SEC)
    return tbl


def stat_row(s, items, t=2.4, h=1.7):
    n = len(items); gap = 0.3; total = 12.0
    w = (total - gap * (n - 1)) / n
    l = 0.7
    for big, label, sub in items:
        panel(s, l, t, w, h, fill=SURF, line=GOLD)
        _, tf = box(s, l, t + 0.2, w, h - 0.4, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, big, 30, GOLD, bold=True, font=DISP, align=PP_ALIGN.CENTER, space_after=4, first=True)
        para(tf, label.upper(), 10.5, MUTED, bold=True, font=BODY, align=PP_ALIGN.CENTER, space_after=2)
        if sub:
            para(tf, sub, 10, SEC, font=BODY, align=PP_ALIGN.CENTER, space_after=0)
        l += w + gap


# 1 TITLE
s = slide()
_, tf = box(s, 1, 2.3, 11.33, 2.9, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "ELEVATE HOMES", 60, GOLD2, bold=True, font=DISP, align=PP_ALIGN.CENTER, space_after=10, first=True)
para(tf, "DUBAI RESALE INTELLIGENCE  ·  BROKERAGE + PRISM SAAS", 15, MUTED, bold=True, font=BODY, align=PP_ALIGN.CENTER, space_after=14)
para(tf, "End-to-End Investor & Feasibility Briefing · June 2026", 15, SEC, font=BODY, align=PP_ALIGN.CENTER, space_after=6)
para(tf, "Sambhav Gupta — 3.5 yrs Dubai real estate · Senior Investment Analyst, Sobha Realty", 12, MUTED, font=BODY, align=PP_ALIGN.CENTER)
notes(s, "One-line: We are not another brokerage — we are a data-led resale engine that uses DLD pricing intelligence to win exclusive mandates and close with higher trust, and we license the same engine as SaaS.")

# 2 THE MOMENT
s = slide()
kicker_title(s, "The Context", "The Market Flipped to Resale.\nWe Were Built for This Moment.")
stat_row(s, [
    ("-27%", "Off-plan secondary Q1 2026", "flip strategies broke"),
    ("+48%", "Ready/resale values YoY", "buyers want certainty"),
    ("AED 53B", "Secondary volume Q1 2026", "~14,399 deals"),
])
_, tf = box(s, 0.7, 4.5, 12, 1.2)
para(tf, "~65,000 units delivering through 2026 → a wave of owners who must price to sell. "
         "Sellers need proof, not a sales pitch.", 13, SEC, font=BODY, first=True)
para(tf, "Sources: Dubai Land Department Q1 2026; figures directional, validated live by PRISM.", 10.5, MUTED, font=BODY, space_after=0)
notes(s, "The war/market shock pushed capital from 2-year off-plan bets into immediate-liquidity resale. That is exactly the segment PRISM serves.")

# 3 PROBLEM
s = slide()
kicker_title(s, "The Problem", "Sellers Guess. Buyers Distrust.\nAgents Have No Evidence.")
card(s, 0.7, 2.5, 5.9, 2.0, "For Sellers", ["No transparent DLD-based pricing tool", "Opinion-based agent valuations", "Need a defensible exit price fast"])
card(s, 6.75, 2.5, 5.9, 2.0, "For Buyers", ["Can't verify fair vs overpriced", "Portals show asking, not sold prices", "Need data before committing capital"])
card(s, 0.7, 4.65, 5.9, 2.0, "For Brokerages", ["Portal leads split across 5,700+ agents", "No tool to win mandates with evidence", "Spend buys raw leads, not mandates"])
card(s, 6.75, 4.65, 5.9, 2.0, "The Gap", ["DLD data is public but not digestible", "Nobody closes data→trust→mandate→fee", "The window is now: motivated sellers"])
notes(s, "The exclusive mandate is the highest-value, lowest-competition asset in brokerage. Evidence is how you win it.")

# 4 SOLUTION
s = slide()
kicker_title(s, "The Solution", "Elevate Homes + the PRISM Engine")
_, tf = box(s, 0.7, 1.95, 12, 0.7)
para(tf, "A DLD-powered pricing-intelligence layer that generates trust, captures leads, converts them into "
         "exclusive resale mandates — and licenses the same engine as SaaS.", 13, SEC, font=BODY, first=True)
card(s, 0.7, 2.8, 3.9, 2.4, "PRISM Engine", ["Hybrid valuation, 66,413 DLD rows", "Log-size elasticity 0.2733", "Recency + size-proximity weighting", "Confidence diagnostics"])
card(s, 4.75, 2.8, 3.9, 2.4, "Pricing Proof", ["Shows the actual comparable sales", "Date, project, rooms, AED/sqft", "Not a black-box number", "Evidence sellers take to market"])
card(s, 8.8, 2.8, 3.85, 2.4, "Mandate Engine", ["Valuation -> lead -> call", "-> exclusive mandate", "-> listed -> closed -> commission", "Tracked in the built-in CRM"])
panel(s, 0.7, 5.45, 11.95, 1.2, fill=RGBColor(0x1A, 0x16, 0x0C), line=GOLD)
_, tf = box(s, 0.95, 5.6, 11.4, 0.95, anchor=MSO_ANCHOR.MIDDLE)
para(tf, '"We are not another brokerage. We are a data-led resale engine that uses pricing intelligence to '
         'acquire mandates and close transactions with higher trust."', 16, GOLD2, italic=True, font=DISP, first=True)
notes(s, "Lead every investor conversation with that line.")

# 5 PRODUCT (LIVE)
s = slide()
kicker_title(s, "The Product · Working Prototype", "Three Flows. One Engine. Live Today.")
card(s, 0.7, 2.4, 3.9, 2.9, "Seller · /sell",
     ["1. Area, project, beds, size", "2. Fair-value band + liquidity", "3. Comparable DLD rows shown",
      "4. Lead captured -> CRM", "5. Book pricing strategy call", "6. -> Exclusive mandate"])
card(s, 4.75, 2.4, 3.9, 2.9, "Buyer · /deal-check",
     ["1. Listing + asking price", "2. Verdict: fair / over / under", "3. Downloadable Pricing Proof",
      "4. Lead captured -> CRM", "5. Let us negotiate this unit", "6. -> Buyer representation"])
card(s, 8.8, 2.4, 3.85, 2.9, "Brokers · /brokers + /crm",
     ["1. Broker early-access signup", "2. Agency + RERA BRN", "3. Feeds SaaS pipeline",
      "4. CRM: New -> Mandate -> Closed", "5. Notes, status, history", "6. -> Network + subscriptions"])
panel(s, 0.7, 5.5, 11.95, 1.1, fill=RGBColor(0x10, 0x1A, 0x14), line=GREEN)
_, tf = box(s, 0.95, 5.62, 11.4, 0.9, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "LIVE NOW: 66,413 clean DLD records · PRISM Hybrid v1 · /healthz ok · /api/engine -> "
         "\"prism-hybrid-v1\" · seller + buyer + broker funnels + branded reports + agent CRM all working.",
     12, SEC, font=BODY, first=True)
notes(s, "This is the differentiator vs a slideware startup: the product already runs.")

# 5B FEATURE LIST
s = slide()
kicker_title(s, "Feature List · No Ambiguity", "Exactly What We're Building")
rows = [
    ["Module", "What it does - precisely", "Status"],
    ["PRISM hybrid engine", "DLD comp median + log-size 0.2733 + recency/size weighting + diagnostics", "LIVE"],
    ["Seller valuation /sell", "Fair value, 3-tier listing range, liquidity, target-price check", "LIVE"],
    ["Buyer deal-check", "Verdict (Buy/Fair/Over) + fair value + P25-P75 band position", "LIVE"],
    ["Comparable transactions", "Actual DLD rows inline: date, project, rooms, sqft, price, AED/sqft", "LIVE"],
    ["Pricing Proof report /report", "Branded print-to-PDF with comparables + methodology", "LIVE"],
    ["Lead capture /api/leads", "Soft-gate; stores intent + property + valuation snapshot", "LIVE"],
    ["Agent CRM /crm", "Pipeline New->Mandate->Closed, notes, history, stage control", "LIVE"],
    ["Broker signup /brokers", "Agency + RERA BRN early access -> CRM", "LIVE"],
    ["Engine API /api/engine", "Programmatic engine + health; foundation for SaaS/licensing", "LIVE"],
]
table(s, 0.7, 1.95, 11.95, rows, [3.0, 7.35, 1.6], row_h=0.4, fsize=9.5)
panel(s, 0.7, 6.15, 11.95, 0.95, fill=RGBColor(0x1A, 0x16, 0x0C), line=GOLD)
_, tf = box(s, 0.95, 6.27, 11.45, 0.75, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Planned (Phase 2-3 - explicitly NOT built yet): SaaS subscriptions & billing (AED 199-2,500/mo) · "
         "pay-per-report credits · buyer-requirements auto-matching · white-label / enterprise data API · "
         "full 7.71 GB DLD cache in production (builder ready).", 10.5, SEC, font=BODY, first=True, line=1.1)
notes(s, "Unambiguous: green = working in the prototype today; the planned list is what we explicitly have NOT built yet.")

# 6 COMPETITIVE
s = slide()
kicker_title(s, "Competitive Moat", "What No One Else Closes")
rows = [
    ["Platform", "What they do", "What's missing", "Loop?"],
    ["DXB Interact", "Free raw DLD data portal", "No journey, no lead, no mandate, no revenue", "No"],
    ["Property Monitor", "B2B AVM for institutions", "Enterprise-only; no consumer; no brokerage", "No"],
    ["Bayut TruEstimate", "300K+ AI valuations -> 5,700 agents", "Distributes trust; no exclusivity; no comps shown", "No"],
    ["Dubai REST (DLD)", "Official govt AVM", "Generic; no comparables; no actionability", "No"],
    ["Elevate Homes", "Data->Proof->Lead->Mandate->Closed->Fee + SaaS", "Nothing - we close the full loop", "YES"],
]
table(s, 0.7, 2.4, 11.95, rows, [2.3, 3.9, 4.55, 1.2], row_h=0.62, fsize=10.5)
_, tf = box(s, 0.7, 6.4, 12, 0.6)
para(tf, "The moat compounds: every closed mandate adds a private comp the incumbents never see.", 11.5, MUTED, font=BODY, first=True)
notes(s, "Incumbents generate data OR leads OR run a brokerage. None close the loop.")

# 7 BUSINESS MODEL
s = slide()
kicker_title(s, "Business Model · Blended", "Brokerage Funds It. SaaS Scales It.")
rows = [
    ["Revenue stream", "Y1", "Y2", "Y3"],
    ["1 · Resale mandates (GCI)", "1.91M", "5.73M", "12.22M"],
    ["2 · CPL agent pipeline", "0.95M", "2.42M", "4.24M"],
    ["3 · Founder personal deals", "0.64M", "0.81M", "0.81M"],
    ["4 · PRISM SaaS subscriptions", "-", "0.74M", "1.85M"],
    ["5 · Per-report credits / area packs", "0.06M", "0.20M", "0.40M"],
    ["6 · Enterprise data / API licensing", "-", "0.18M", "0.46M"],
    ["COMBINED (AED)", "3.56M", "10.07M", "19.98M"],
]
table(s, 0.7, 2.35, 8.4, rows, [4.5, 1.3, 1.3, 1.3], row_h=0.5, fsize=11)
card(s, 9.4, 2.35, 3.25, 4.0, "Two arms, one engine",
     ["Brokerage = transaction GCI (~98% gross), funds Year 1.",
      "SaaS = recurring + usage; builds the moat & the multiple.",
      "", "SaaS tiers (AED/mo):", "Free (own agents)", "199 Starter · 499 Pro",
      "999 Team · 2,500+ Enterprise"])
_, tf = box(s, 0.7, 6.55, 12, 0.5)
para(tf, "Forward projections (blended financial model). SaaS share grows ~2% -> 11% -> 14% of revenue.", 10.5, MUTED, font=BODY, first=True)
notes(s, "Brokerage-first execution, blended monetization, dual exit. Don't make SaaS load-bearing in Year 1.")

# 8 UNIT ECONOMICS
s = slide()
kicker_title(s, "Unit Economics", "The Numbers Work From Deal One")
stat_row(s, [
    ("18.1x", "LTV : CAC", "target >10x"),
    ("AED 120K", "Avg GCI / deal", "2% on ~AED 6M"),
    ("Month 7", "Breakeven", "brokerage arm"),
    ("AED 20K", "Founder salary / mo", "only fixed salary"),
], t=2.3, h=1.6)
_, tf = box(s, 0.7, 4.2, 12, 2.6)
para(tf, "Funnel (working assumptions to validate in Sprint 1):", 13, GOLD2, bold=True, font=BODY, first=True, space_after=8)
para(tf, "▸  App submission → qualified lead:  37%", 13, SEC, font=BODY, space_after=6)
para(tf, "▸  Qualified lead → consultation:  28%", 13, SEC, font=BODY, space_after=6)
para(tf, "▸  Consultation → exclusive mandate:  45%", 13, SEC, font=BODY, space_after=10)
para(tf, "Break-even ≈ 1.5 mandates/month — reachable with the founder closing alone.", 12, MUTED, font=BODY, space_after=0)
notes(s, "Be honest: conversion % are assumptions. The model survives worse numbers because the fixed base is one salary.")

# 9 FINANCIALS
s = slide()
kicker_title(s, "Financial Model · 3 Years", "Capital-Light. Cash-Positive Fast.")
stat_row(s, [
    ("AED 3.56M", "Year 1 revenue", "EBITDA 1.75M · 49%"),
    ("AED 10.07M", "Year 2 revenue", "EBITDA 6.53M · 65%"),
    ("AED 19.98M", "Year 3 revenue", "EBITDA 15.12M · 76%"),
], t=2.3, h=1.7)
card(s, 0.7, 4.4, 3.9, 2.2, "Two arms", ["Brokerage = GCI (~98% gross).", "SaaS COGS only AED 120K -> 264K."])
card(s, 4.75, 4.4, 3.9, 2.2, "Breakeven", ["Brokerage ~M7 (GCI carries Y1).", "SaaS self-sustaining ~M18.", "Positive EBITDA from Year 2."])
card(s, 8.8, 4.4, 3.85, 2.2, "Indicative Y3 value", ["SaaS ~1.85M ARR x ~10x", "+ brokerage EBITDA x ~4x", "-> ~AED 40-60M · dual exit"])
notes(s, "Forward projections; multiples illustrative, not a quote.")

# 10 SCALABILITY I
s = slide()
kicker_title(s, "Scalability · The Flywheel", "A Brokerage That Gets Smarter\nWith Every Deal")
_, tf = box(s, 0.7, 2.5, 12, 0.7)
para(tf, "More valuations → more mandates & closings → private comps PRISM owns → smarter engine → "
         "sharper pricing → more trust → more valuations.", 13.5, GOLD2, font=BODY, first=True, line=1.2)
card(s, 0.7, 3.6, 3.9, 2.6, "Zero marginal data cost",
     ["DLD refreshes nightly.", "Zero-dependency Node engine.", "Next area / next 100K rows = compute, not headcount."])
card(s, 4.75, 3.6, 3.9, 2.6, "Product-led, not payroll-led",
     ["Mandates from the funnel, not 50 cold-callers.", "Hire to proven demand:", "1 agent / 8-10 active mandates."])
card(s, 8.8, 3.6, 3.85, 2.6, "Same engine, two buyers",
     ["In-house agents use PRISM free.", "External agents/devs pay.", "Every subscriber surfaces a deal we can capture."])
notes(s, "This is why it's PropTech, not 'another brokerage' — the asset compounds.")

# 11 SCALABILITY II
s = slide()
kicker_title(s, "Scalability · Expansion Path", "Win Density First, Then Widen")
items = [
    ("Stage 1 — Depth (Now -> M6)", "Dominate 3-5 high-liquidity communities (Business Bay, Marina, JVC, Downtown, Dubai Hills). Repeatable mandates + a private comp set no portal has."),
    ("Stage 2 — Breadth (M6 -> M18)", "Roll the same playbook across all Dubai resale communities. Activate the broker network + launch PRISM SaaS to external agents."),
    ("Stage 3 — Platform (M18 -> Y3)", "Enterprise data/API + white-label PRISM for brokerages & banks. Layer buyer-side + off-plan referral channels on the same lead engine."),
    ("Stage 4 — Geographic (Y3+)", "Same architecture, new registries: Abu Dhabi, then GCC markets with open transaction data. The engine is portable; only the data source changes."),
]
t = 2.45
for head, body in items:
    panel(s, 0.7, t, 11.95, 1.0, fill=SURF, line=GOLD)
    _, tf = box(s, 0.95, t + 0.12, 11.5, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, head, 13, GOLD2, bold=True, font=BODY, first=True, space_after=3)
    para(tf, body, 11, SEC, font=BODY, space_after=0, line=1.05)
    t += 1.12
notes(s, "Each stage is gated by proof from the previous one — capital follows traction.")

# 12 ROADMAP
s = slide()
kicker_title(s, "Roadmap · First 90 Days", "Traction Now. Scale on Proof.")
items = [
    ("NOW · Days 1-30", "Resale Proof MVP (prototype built): seller + buyer funnels, lead gate, Pricing Proof PDF, agent CRM, 10 area pages, 3 seller campaigns.", "KPIs: 300 visitors · 50 submissions · 5 seller calls · 2 mandates"),
    ("NEXT · Days 31-60", "Mandate Acquisition: full 7.71 GB DLD canonical cache, report share links, agent dashboard v1, retargeting, activate 2 resale agents on proof.", "KPIs: 100+ submissions/mo · 10+ calls · 5+ mandates · 1-2 closings"),
    ("SCALE · Days 61-90", "Systemize & open SaaS beta: area heatmap, seller dashboard, gated hiring, PRISM SaaS beta to first external brokers.", "KPIs: 300+ submissions/mo · 15+ mandates · 3-5 closings · first SaaS signups"),
]
t = 2.5
for head, body, kpi in items:
    panel(s, 0.7, t, 11.95, 1.25, fill=SURF, line=GOLD)
    _, tf = box(s, 0.95, t + 0.12, 11.5, 1.05, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, head, 12.5, GOLD2, bold=True, font=BODY, first=True, space_after=3)
    para(tf, body, 11, SEC, font=BODY, space_after=3, line=1.04)
    para(tf, kpi, 10.5, GOLD, font=BODY, space_after=0)
    t += 1.4
notes(s, "Show traction now; scale only on proof.")

# 13 FEASIBILITY I
s = slide()
kicker_title(s, "Feasibility Check · 1 of 2", "Does This Actually Work?")
def feas_card(l, t, dim, rating, color, why):
    panel(s, l, t, 3.9, 1.85, fill=SURF, line=color)
    _, tf = box(s, l + 0.22, t + 0.16, 3.5, 1.6)
    para(tf, dim + "   [" + rating + "]", 13, color, bold=True, font=BODY, first=True, space_after=5)
    para(tf, why, 10.5, SEC, font=BODY, space_after=0, line=1.05)
feas_card(0.7, 2.35, "Market", "HIGH", GREEN, "AED 53B secondary/qtr; motivated sellers; structural shift to resale. Demand proven.")
feas_card(4.75, 2.35, "Technical", "HIGH", GREEN, "Prototype live: engine, funnels, CRM, reports. Zero-dep, nightly refresh.")
feas_card(8.8, 2.35, "Financial", "MED-HIGH", GOLD, "Capital-light; GCI funds Y1; breakeven ~M7. Sensitive to CAC + conversion.")
feas_card(0.7, 4.4, "Operational", "MEDIUM", GOLD, "Founder-led, commission-only agents, gated hiring. 5-min response SLA.")
feas_card(4.75, 4.4, "Regulatory", "MEDIUM", GOLD, "RERA licence + BRN; DLD redistribution terms. Known and navigable.")
feas_card(8.8, 4.4, "Competitive", "MED-HIGH", GOLD, "No incumbent closes the loop; moat = private mandate comps. Watch the portals.")
notes(s, "Honest multi-dimensional read from an operator's seat.")

# 14 FEASIBILITY II
s = slide()
kicker_title(s, "Feasibility Check · 2 of 2", "Top Risks — and How We Kill Them")
rows = [
    ["Risk", "Mitigation"],
    ["Full DLD access (Pulse pending)", "66K cached works now; canonical builder ready; nightly refresh live"],
    ["Mandate conversion below model", "Founder closes first; don't hire until proven; 60-90 day exclusives"],
    ["CAC inflation on portals", "Own-channel SEO/content/WhatsApp; NRI targeting; report = marketing"],
    ["Brokerage shakeout", "Lean by design: one salary, commission-only. Built to be the survivor"],
    ["Incumbent copies comparables", "Speed + the mandate-comp flywheel + consumer->brokerage loop"],
]
table(s, 0.7, 2.35, 11.95, rows, [4.0, 7.95], row_h=0.56, fsize=10.5)
panel(s, 0.7, 5.95, 11.95, 1.05, fill=RGBColor(0x10, 0x1A, 0x14), line=GREEN)
_, tf = box(s, 0.95, 6.05, 11.4, 0.9, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "VERDICT: GO (conditional). Capital-light, product already built, demand validated, honest moat. "
         "Conditions: prove mandate conversion in 1-2 areas before scaling spend, and secure the full DLD path. "
         "Risk is in execution, not market or technology.", 12, GOLD2, font=BODY, first=True, line=1.1)
notes(s, "Kill criteria: CAC>AED25K/mandate, conversion<20% with founder closing, or DLD access foreclosed.")

# 15 TEAM
s = slide()
kicker_title(s, "Team & Credibility", "Domain First. Code Followed.")
card(s, 0.7, 2.5, 3.9, 2.6, "The operator", ["Sambhav Gupta", "3.5 yrs Dubai real estate", "Senior Investment Analyst, Sobha Realty", "Lives the seller/buyer/agent problem daily"])
card(s, 4.75, 2.5, 3.9, 2.6, "The research", ["MBA hedonic-pricing dissertation", "DLD branded-residence premiums", "Size elasticity ~0.27 -> engine's 0.2733", "PRISM re-tests it on live data"])
card(s, 8.8, 2.5, 3.85, 2.6, "The build", ["Working prototype shipped solo", "AI-assisted development", "Bottleneck is domain depth, not a dev team", "We have the domain"])
panel(s, 0.7, 5.35, 11.95, 1.1, fill=RGBColor(0x1A, 0x16, 0x0C), line=GOLD)
_, tf = box(s, 0.95, 5.5, 11.4, 0.85, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Hiring philosophy: earn the right to hire. One founder salary today; commission-only agents added "
         "only as proven mandate flow demands — gated, not speculative.", 12.5, SEC, font=BODY, first=True)
notes(s, "Founder-market fit is the story: a domain expert who shipped the tool.")

# 16 THE ASK
s = slide()
kicker_title(s, "The Ask", "What We Need to Win")
stat_row(s, [
    ("AED 630K", "Seed · 4 tranches", "milestone-gated"),
    ("20%", "Equity", ""),
    ("Month 7", "Breakeven", "~AED 5M cash Y3"),
], t=2.25, h=1.55)
card(s, 0.7, 4.1, 5.9, 2.4, "From the Incubator",
     ["PropTech + brokerage mentorship", "Pilot users + test mandates via network", "RERA licensing + setup guidance", "Investor-readiness review"])
card(s, 6.75, 4.1, 5.9, 2.4, "Use of Funds",
     ["RERA licence + office setup", "Mandate-acquisition marketing (gated)", "Full DLD pipeline + SaaS build", "First agents — only on proven flow"])
notes(s, "Closing line: every competitor generates data OR leads OR runs a brokerage. Elevate Homes closes the full loop — and licenses the engine.")

# 17 CLOSING
s = slide()
_, tf = box(s, 1, 1.7, 11.33, 2.0, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "ELEVATE HOMES", 52, GOLD2, bold=True, font=DISP, align=PP_ALIGN.CENTER, space_after=8, first=True)
para(tf, "RESALE INTELLIGENCE · BROKERAGE + PRISM SAAS", 14, MUTED, bold=True, font=BODY, align=PP_ALIGN.CENTER)
stat_row(s, [
    ("66,413", "DLD records live", ""),
    ("AED 20M", "Y3 revenue", "modeled"),
    ("~AED 40-60M", "Y3 valuation", "indicative"),
    ("GO", "Feasibility verdict", "conditional"),
], t=4.0, h=1.6)
_, tf = box(s, 1, 6.5, 11.33, 0.6)
para(tf, "Sambhav Gupta · Elevate Homes · June 2026 · prototype live", 12, MUTED, font=BODY, align=PP_ALIGN.CENTER, first=True)
notes(s, "Thank you. Prototype is live and clickable — happy to demo /sell, /deal-check, /crm, /report, /deck.")

out = os.path.join(os.path.dirname(__file__), "Elevate-Homes-End-to-End-Deck.pptx")
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
